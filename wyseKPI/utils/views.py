from django.shortcuts import render
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE, MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
import numpy as np
from PIL import Image, ImageDraw
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.xmlchemy import OxmlElement
from rest_framework import status
from django.http import Http404
from rest_framework.views import APIView
from rest_framework.response import Response
from .models import AggregateRating, Workforce, KPIProject, KPI, KPIRating

# Create your views here.

class MultipleEmployee(APIView):
    def getEmployeeAvgRating(self, project_id, employee_id):
        try:
            employeeData = AggregateRating.objects.filter(project_id_id = project_id, emp_id_id=employee_id, review_type='review')
            return employeeData
        except:
            raise Http404

    def getEmployeeDetails(self, project_id):
        employeeList = []
        try:
            employee = Workforce.objects.filter(project_id_id=project_id)
            for i in employee:
                employeeList.append(i.emp_id_id)
            return employeeList
        except:
            raise Http404

    def getEmployeeKpiRating(self, project_id, emp_id, kpi_id):
        kpiData = []
        try:
            rating = KPIRating.objects.filter(project_id_id=project_id, emp_id_id=emp_id, kpi_id_id=kpi_id, review_type='review')
            for i in rating:
                kpiData.append(i.rating)
            return kpiData
        except:
            raise Http404

    def getKpiList(self, project_id):
        kpiList = []
        try:
            kpis = KPIProject.objects.filter(project_id_id=project_id)
            for i in kpis:
                kpiList.append(i.kpi_id_id)
            return kpiList
        except:
            raise Http404

    def getKpiName(self, kpi_id):
        try:
            name = KPI.objects.get(kpi_id=kpi_id)
            return name.kpi_name
        except:
            raise Http404

    def SubElement(self,parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

    def image(self, path):
        img = Image.open(path).convert("RGB")
        npImage = np.array(img)
        h, w = img.size

        # Create same size alpha layer with circle
        alpha = Image.new('L', img.size, 0)
        draw = ImageDraw.Draw(alpha)
        draw.pieslice([0, 0, h, w], 0, 360, fill=255)

        # Convert alpha Image to numpy array
        npAlpha = np.array(alpha)

        # Add alpha layer to RGB
        npImage = np.dstack((npImage, npAlpha))

        # Save with alpha
        return Image.fromarray(npImage).save("download.png")

    def makeParaBulletPointed(self, para):
        """Bullets are set to Arial,
            actual text can be a different font"""
        pPr = para._p.get_or_add_pPr()
        ## Set marL and indent attributes
        pPr.set('marL', '171450')
        pPr.set('indent', '171450')
        ## Add buFont
        _ = self.SubElement(parent=pPr,
                       tagname="a:buFont",
                       typeface="Arial",
                       panose="020B0604020202020204",
                       pitchFamily="34",
                       charset="0"
                       )
        ## Add buChar
        _ = self.SubElement(parent=pPr,
                       tagname='a:buChar',
                       char="-")

    ## <<< ------------------------ All Page Header ------------------------------->>>
    def allPageHeader(self, slide, headerText):
        mainLeft = Inches(0.2)
        mainTop = Inches(0)
        mainWidth = Inches(9)
        mainHeight = Inches(0.6)
        mainTx_box = slide.shapes.add_textbox(mainLeft, mainTop, mainWidth, mainHeight)
        tf = mainTx_box.text_frame
        p = tf.add_paragraph()
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p.text = headerText
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(18)

    ## <<< ---------------------------- Adding Image ----------------- >>>
    def addImage(self, slide):
        sideTextBoxLeft = Inches(0)
        sideTextBoxTop = Inches(.7)
        sideTextBoxWidth = Inches(2)
        sideTextBoxHeight = Inches(6.8)
        sideTextBox = slide.shapes.add_textbox(sideTextBoxLeft, sideTextBoxTop, sideTextBoxWidth, sideTextBoxHeight)
        left = Inches(0.3)
        top = Inches(3)
        width = height = Inches(1.5)
        self.image("download.jpg")
        slide.shapes.add_picture("download.png", left, top, width, height)
        # tf = sideTextBox.text_frame
        # p=tf.add_paragraph()

    ## <<< ---------------------------- Adding Image ----------------- >>>
    def addTextBox(self, slide, left, top, width, height):
        textBox = slide.shapes.add_textbox(left, top, width, height)
        return textBox

    ## <<<---------------------- Method for  ticket block generation ------------------------------->>>
    def tickectBlock(self, slide, verLenGreen, verStartGreen, firstTextBoxLeft, firstTextBoxTop, firstTextBoxWidth,
                     firstTextBoxHeight, secondTextBoxHeight, lineVerleft=2.6, lineHorleft=2.8):
        # Creating Vertical Line
        line1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LINE_INVERSE, Inches(lineVerleft), Inches(verStartGreen),
                                       Inches(0),
                                       Inches(verLenGreen))
        line1.line.color.rgb = RGBColor(26, 201, 101)
        verStartGreen = verStartGreen + verLenGreen

        # adding Diamond
        d = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.DIAMOND, Inches(firstTextBoxLeft - 0.298), Inches(firstTextBoxTop + 0.07), Inches(0.2),
            Inches(0.2)
        )
        fill = d.fill
        fill.solid()
        fore_color = fill.fore_color
        fore_color.rgb = RGBColor(26, 201, 101)

        # Adding Ticket Box
        textBox = self.addTextBox(slide, Inches(firstTextBoxLeft), Inches(firstTextBoxTop), Inches(firstTextBoxWidth),
                             Inches(firstTextBoxHeight))

        # Adding Story Point Box
        textBox2 = self.addTextBox(slide, Inches(firstTextBoxLeft), Inches(firstTextBoxTop + .2), Inches(firstTextBoxWidth),
                              Inches(secondTextBoxHeight))
        # Inserting Value to Tickect Box
        self.addParagraph(textBox, ticketNumber[i], rgb=(0, 0, 0), fontsize=20, bold=True, para=PP_PARAGRAPH_ALIGNMENT.LEFT)

        # Inserting Value to Story Point Box
        self.addParagraph(textBox2, 'Start Date: 21/11/2022 | End Date: -', fontsize=14, bold=True, rgb=(0, 0, 0),
                     para=PP_PARAGRAPH_ALIGNMENT.LEFT)
        self.addParagraph(textBox2, 'Story Pt: 3-5 | Status: Completed', bold=True, rgb=(0, 0, 0),
                     para=PP_PARAGRAPH_ALIGNMENT.LEFT)
        # Adding Horizontal Line
        line2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LINE_INVERSE, Inches(lineHorleft),
                                       Inches(secondTextBoxHeight + firstTextBoxTop + .3), Inches(4),
                                       Inches(0))
        line2.line.color.rgb = RGBColor(191, 191, 191)
        return verStartGreen

    ## <<< ---------------------------- Fill color in the Box ----------------- >>>
    def fillColorInBox(self, slide, value=(26, 201, 101)):
        fill = slide.fill
        fill.solid()
        fore_color = fill.fore_color
        fore_color.rgb = RGBColor(value[0], value[1], value[2])

    ## <<< ---------------------------- Add Paragraph ----------------- >>>
    def addParagraph(self, slide, textValue, bold=False, fontsize=12, lineSpacing=15, rgb=(255, 255, 255),
                     para=PP_PARAGRAPH_ALIGNMENT.CENTER, verPara=MSO_VERTICAL_ANCHOR.BOTTOM, val=False):
        tf = slide.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = textValue
        p.alignment = para
        tf.vertical_anchor = verPara
        p.font.bold = bold
        p.line_spacing = Pt(lineSpacing)
        p.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
        p.font.size = Pt(fontsize)
        if val:
            self.makeParaBulletPointed(p)

    def get(self, request, project_id):
        # <<< ---------------------- Title Page --------------------------- >>>
        prs = Presentation("slide.pptx")
        slide = prs.slides[0]
        mainDataText = "DE Team - November Report"
        for shape in slide.shapes:
            if shape.shape_type == 17:
                slide.shapes.element.remove(shape.element)
        tx_box = self.addTextBox(slide, left=Inches(6.6), top=Inches(5.8), width=Inches(6.5), height=Inches(1))
        self.addParagraph(tx_box, mainDataText, True, 30, verPara=MSO_VERTICAL_ANCHOR.TOP)

        # <<< ---------------------- Page 2 Content --------------------------- >>>
        slide1 = prs.slides[1]
        contentText = ["November Deliverables", "December Plans", "Team Performance Report",
                       "Project Task Deliverables (NOV)"]
        txBox = self.addTextBox(slide1, left=Inches(0.5), top=Inches(1.1), width=Inches(7), height=Inches(6.3))
        for i in contentText:
            p = self.addParagraph(txBox, i, False, 20, 30, (0, 0, 0), para=PP_PARAGRAPH_ALIGNMENT.LEFT, verPara=MSO_VERTICAL_ANCHOR.TOP, val=True)

        employeeRating = []
        employeeList = self.getEmployeeDetails(project_id)
        for i in range(len(employeeList)):
            kpi = []
            employeeRating.append({employeeList[i]:[]})
            employeeAvgData = self.getEmployeeAvgRating(project_id, employeeList[i])
            for j in employeeAvgData:
                employeeRating[i][employeeList[i]].append(j.rating)
            kpiList = self.getKpiList(project_id)
            for k in range(len(kpiList)):
                kpi.append({'kpi Id': kpiList[k], 'kpi_name':''})
                kpiName = self.getKpiName(kpiList[k])
                kpi[k]['kpi_name']=kpiName
                kpi[k]['KPI Value'] = self.getEmployeeKpiRating(project_id, employeeList[i], kpiList[k])
            # <<<<<< -------------- Employee  Page ---------------------------- >>>>>
            # <------------------- page 3 -------------->
            # slideMain = prs.slide_layouts[1]
            # slide = prs.slides.add_slide(slideMain)
            # headerText = "2022  |  November – Deliverables & December Planned Task"
            # self.allPageHeader(slide, headerText)
            # self.addImage(slide)
            # sideBoxName = self.addTextBox(slide, Inches(0.15), Inches(4.5), Inches(1.8), Inches(0.4))
            # self.addParagraph(sideBoxName, "Ranjeet Verma", rgb=(70, 130, 180), fontsize=16, bold=True)
            # ticket = [[1, 2, 7, 78], [1, 2]]
            #
            # ### <<< ---------------------- Static values ---------------------->>>
            # hor = 2.5
            # ver = 1
            # firstTextBoxLeft = 2.8
            # firstTextBoxTop = 1.6
            # firstTextBoxWidth = 4
            # firstTextBoxHeight = 0.4
            # secondTextBoxHeight = 0.8
            # verLenGreen = 1.3
            # verStartGreen = 1.4
            # monthTextValue = "NOV 22"
            # ticketNumber = ['DE-824', 'DE-825', 'DE-826', 'DE-827', 'DE-828', 'DE-829', 'DE-830', 'DE-831', 'DE-832',
            #                 'DE-833', 'DE-834', 'DE-835', 'DE-836', 'DE-837', 'DE-838', 'DE-839']
            #
            # ### <<< --------------------- slide content Loop --------------------- >>>
            # for i in range(len(ticket[0]) + len(ticket[1])):
            #     if i > 7:
            #         if i < 12:
            #             if i == 8:
            #                 slide = prs.slides.add_slide(slideMain)
            #                 header = "2022  |  November - Deliverables & December Planned Task"
            #                 self.allPageHeader(slide, header)
            #                 self.addImage(slide)
            #                 sideBoxName = self.addTextBox(slide, Inches(0.15), Inches(4.5), Inches(1.8), Inches(0.4))
            #                 self.addParagraph(sideBoxName, "Ranjeet Verma", rgb=(70, 130, 180), fontsize=16, bold=True)
            #                 monthBox = self.addTextBox(slide, Inches(hor), Inches(ver), Inches(1.5), Inches(0.4))
            #                 self.addParagraph(monthBox, monthTextValue, fontsize=18, bold=True)
            #                 self.fillColorInBox(monthBox)
            #                 verStartGreen = self.tickectBlock(slide, verLenGreen, verStartGreen, firstTextBoxLeft,
            #                                              firstTextBoxTop,
            #                                              firstTextBoxWidth, firstTextBoxHeight, secondTextBoxHeight)
            #                 firstTextBoxTop += 1.3
            #             elif i < len(ticket[0]):
            #                 verStartGreen = self.tickectBlock(slide, verLenGreen, verStartGreen, firstTextBoxLeft,
            #                                              firstTextBoxTop,
            #                                              firstTextBoxWidth, firstTextBoxHeight, secondTextBoxHeight)
            #                 firstTextBoxTop += 1.3
            #             elif i == len(ticket[0]):
            #                 monthTextValue = "DEC 22"
            #                 monthBox = self.addTextBox(slide, Inches(hor), Inches(firstTextBoxTop + 0.2), Inches(1.5),
            #                                       Inches(0.4))
            #                 self.fillColorInBox(monthBox, value=(240, 228, 60))
            #                 firstTextBoxTop += 0.8
            #                 self.addParagraph(monthBox, monthTextValue, fontsize=18, bold=True, rgb=(0, 0, 0))
            #                 verStartGreen = firstTextBoxTop - 0.2
            #                 verStartGreen = self.tickectBlock(slide, verLenGreen, verStartGreen, firstTextBoxLeft,
            #                                              firstTextBoxTop,
            #                                              firstTextBoxWidth, firstTextBoxHeight, secondTextBoxHeight)
            #                 firstTextBoxTop += 1.3
            #             else:
            #                 verStartGreen = self.tickectBlock(slide, verLenGreen, verStartGreen, firstTextBoxLeft,
            #                                              firstTextBoxTop,
            #                                              firstTextBoxWidth, firstTextBoxHeight, secondTextBoxHeight)
            #                 firstTextBoxTop += 1.3
            #             if i == 11:
            #                 firstTextBoxLeft = 7.4
            #                 firstTextBoxTop = 1.6
            #                 verStartGreen = 1.4
            #                 lineVerleft = 7.2
            #         else:
            #             if i < len(ticket[0]):
            #                 verStartGreen = self.tickectBlock(slide, verLenGreen, verStartGreen, firstTextBoxLeft,
            #                                              firstTextBoxTop,
            #                                              firstTextBoxWidth, firstTextBoxHeight, secondTextBoxHeight,
            #                                              lineVerleft=7.2, lineHorleft=7.4)
            #                 firstTextBoxTop += 1.3
            #             elif i == len(ticket[0]):
            #                 monthTextValue = "DEC 22"
            #                 monthBox = self.addTextBox(slide, Inches(7.1), Inches(firstTextBoxTop + 0.2), Inches(1.5),
            #                                       Inches(0.4))
            #                 self.fillColorInBox(monthBox, value=(240, 228, 60))
            #                 firstTextBoxTop += 0.8
            #                 self.addParagraph(monthBox, monthTextValue, fontsize=18, bold=True, rgb=(0, 0, 0))
            #                 verStartGreen = firstTextBoxTop - 0.2
            #                 verStartGreen = self.tickectBlock(slide, verLenGreen, verStartGreen, firstTextBoxLeft,
            #                                              firstTextBoxTop,
            #                                              firstTextBoxWidth, firstTextBoxHeight, secondTextBoxHeight,
            #                                              lineVerleft=7.2, lineHorleft=7.4)
            #                 firstTextBoxTop += 1.3
            #             else:
            #                 verStartGreen = self.tickectBlock(slide, verLenGreen, verStartGreen, firstTextBoxLeft,
            #                                              firstTextBoxTop,
            #                                              firstTextBoxWidth, firstTextBoxHeight, secondTextBoxHeight,
            #                                              lineVerleft=7.2, lineHorleft=7.4)
            #                 firstTextBoxTop += 1.3
            #             if i == 7:
            #                 hor = 2.5
            #                 ver = 1
            #                 firstTextBoxLeft = 2.8
            #                 firstTextBoxTop = 1.6
            #                 firstTextBoxWidth = 4
            #                 firstTextBoxHeight = 0.4
            #                 secondTextBoxHeight = 0.8
            #                 verLenGreen = 1.3
            #                 verStartGreen = 1.4
            #     elif i == 0:
            #         monthBox = self.addTextBox(slide, Inches(hor), Inches(ver), Inches(1.5), Inches(0.4))
            #         self.fillColorInBox(monthBox)
            #         self.addParagraph(monthBox, monthTextValue, fontsize=18, bold=True)
            #         verStartGreen = self.tickectBlock(slide, verLenGreen, verStartGreen, firstTextBoxLeft, firstTextBoxTop,
            #                                      firstTextBoxWidth, firstTextBoxHeight, secondTextBoxHeight)
            #         firstTextBoxTop += 1.3
            #     else:
            #         if i < 4:
            #             if i < len(ticket[0]):
            #                 verStartGreen = self.tickectBlock(slide, verLenGreen, verStartGreen, firstTextBoxLeft,
            #                                              firstTextBoxTop,
            #                                              firstTextBoxWidth, firstTextBoxHeight, secondTextBoxHeight)
            #                 firstTextBoxTop += 1.3
            #             elif i == len(ticket[0]):
            #                 monthTextValue = "DEC 22"
            #                 monthBox = self.addTextBox(slide, Inches(hor), Inches(firstTextBoxTop + 0.2), Inches(1.5),
            #                                       Inches(0.4))
            #                 self.fillColorInBox(monthBox, value=(240, 228, 60))
            #                 firstTextBoxTop += 0.8
            #                 self.addParagraph(monthBox, monthTextValue, rgb=(0, 0, 0), fontsize=18, bold=True)
            #                 verStartGreen = firstTextBoxTop - 0.2
            #                 verStartGreen = self.tickectBlock(slide, verLenGreen, verStartGreen, firstTextBoxLeft,
            #                                              firstTextBoxTop,
            #                                              firstTextBoxWidth, firstTextBoxHeight, secondTextBoxHeight)
            #                 firstTextBoxTop += 1.3
            #             else:
            #                 verStartGreen = self.tickectBlock(slide, verLenGreen, verStartGreen, firstTextBoxLeft,
            #                                              firstTextBoxTop,
            #                                              firstTextBoxWidth, firstTextBoxHeight, secondTextBoxHeight)
            #                 firstTextBoxTop += 1.3
            #             if i == 3:
            #                 firstTextBoxLeft = 7.4
            #                 firstTextBoxTop = 1.6
            #                 verStartGreen = 1.4
            #                 lineVerleft = 7.2
            #         else:
            #             if i < len(ticket[0]):
            #                 verStartGreen = self.tickectBlock(slide, verLenGreen, verStartGreen, firstTextBoxLeft,
            #                                              firstTextBoxTop,
            #                                              firstTextBoxWidth, firstTextBoxHeight, secondTextBoxHeight,
            #                                              lineVerleft=7.2, lineHorleft=7.4)
            #                 firstTextBoxTop += 1.3
            #             elif i == len(ticket[0]):
            #                 monthTextValue = "DEC 22"
            #                 monthBox = self.addTextBox(slide, Inches(7.1), Inches(firstTextBoxTop + 0.2), Inches(1.5),
            #                                       Inches(0.4))
            #                 self.fillColorInBox(monthBox, value=(240, 228, 60))
            #                 firstTextBoxTop += 0.8
            #                 self.addParagraph(monthBox, monthTextValue, fontsize=18, bold=True, rgb=(0, 0, 0))
            #                 verStartGreen = firstTextBoxTop - 0.2
            #                 verStartGreen = self.tickectBlock(slide, verLenGreen, verStartGreen, firstTextBoxLeft,
            #                                              firstTextBoxTop,
            #                                              firstTextBoxWidth, firstTextBoxHeight, secondTextBoxHeight,
            #                                              lineVerleft=7.2, lineHorleft=7.4)
            #                 firstTextBoxTop += 1.3
            #             else:
            #                 verStartGreen = self.tickectBlock(slide, verLenGreen, verStartGreen, firstTextBoxLeft,
            #                                              firstTextBoxTop,
            #                                              firstTextBoxWidth, firstTextBoxHeight, secondTextBoxHeight,
            #                                              lineVerleft=7.2, lineHorleft=7.4)
            #                 firstTextBoxTop += 1.3
            #             if i == 7:
            #                 hor = 2.5
            #                 ver = 1
            #                 firstTextBoxLeft = 2.8
            #                 firstTextBoxTop = 1.6
            #                 verStartGreen = 1.4
            # <<<------------------------ graph slide ----------------------->>>

            # creating Graphs on PPTs
            slide5_register = prs.slide_layouts[1]
            slide5 = prs.slides.add_slide(slide5_register)

            headerText = str(employeeList[i])+" -  November 2022 Performance Report"
            self.allPageHeader(slide5, headerText)
            self.addImage(slide5)
            sideBoxName = self.addTextBox(slide5, Inches(0.15), Inches(4.5), Inches(1.8), Inches(0.4))
            self.addParagraph(sideBoxName, str(employeeList[i]), rgb=(70, 130, 180), fontsize=16, bold=True)

            # build Graph
            graph_info = CategoryChartData()
            graph_info.categories = [i+1 for i in range(len(kpi[0]['KPI Value']))]
            graph_info.add_series(kpi[0]['kpi_name'], kpi[0]['KPI Value'])
            graph_info.add_series(kpi[1]['kpi_name'], kpi[1]['KPI Value'])
            graph_info.add_series(kpi[2]['kpi_name'], kpi[2]['KPI Value'])
            graph_info.add_series(kpi[3]['kpi_name'], kpi[3]['KPI Value'])

            # add Graph to Slide with positioning
            left_graph = Inches(2.5)
            top_graph = Inches(2.5)
            width_graph = Inches(6)
            height_graph = Inches(4.5)
            graph1_frame = slide5.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS,
                                                   left_graph, top_graph, width_graph
                                                   , height_graph, graph_info)

            graph1 = graph1_frame.chart
            graph1.has_legend = True
            graph1.legend.position = XL_LEGEND_POSITION.BOTTOM
            graph1.legend.include_in_layout = False
            graph1.series[0].smooth = True
            graph1.gridlines = False
            graph1.MajorGridlines = False

            # adding title to graph
            graph1.has_title = True
            graph_name = graph1.chart_title.text_frame

            # adding x and y axis names
            category_axis_title = graph1.category_axis.axis_title
            category_axis_title.text_frame.text = "Week"
            value_axis_title = graph1.value_axis.axis_title
            value_axis_title.text_frame.text = "Score"

            Right = Inches(9)
            Top = Inches(3)
            mainWidth = Inches(4)
            mainHeight = Inches(2)
            mainTx_box = slide5.shapes.add_textbox(Right, Top, mainWidth, mainHeight)
            tf = mainTx_box.text_frame
            p = tf.add_paragraph()
            p.text = "Overall Score : 7.7 \nMeeting all deadlines with accuracy \non docgen"
            p.font.bold = True

        # Multiple Graph
        slide6_register = prs.slide_layouts[1]
        slide6 = prs.slides.add_slide(slide6_register)

        self.allPageHeader(slide6, "DE TEAM WEEKLY PERFORMANCE REPORT")

        # build Graph
        graph_info = CategoryChartData()
        graph_info.categories = [i+1 for i in range(len(employeeRating[0][4]))]
        graph_info.add_series(str(employeeList[0]), employeeRating[0][employeeList[0]])
        graph_info.add_series(str(employeeList[1]), employeeRating[1][employeeList[1]])
        graph_info.add_series(str(employeeList[2]), employeeRating[2][employeeList[2]])
        # graph_info.add_series("Vamshi", (9, 9, 9, 9))
        # graph_info.add_series("Shubhamay", (1, 9, 9, 9))
        # graph_info.add_series("Hemraj", (8, 9, 7, 9))

        # add Graph to Slide with positioning
        left_graph = Inches(6)
        top_graph = Inches(2.3)
        width_graph = Inches(6)
        height_graph = Inches(4.5)
        graph1_frame = slide6.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS,
                                               left_graph, top_graph, width_graph
                                               , height_graph, graph_info)

        graph1 = graph1_frame.chart
        graph1.value_axis.minimum_scale = 0
        graph1.value_axis.maximum_scale = 10
        graph1.has_legend = True
        graph1.legend.position = XL_LEGEND_POSITION.BOTTOM
        graph1.legend.include_in_layout = False
        graph1.series[0].smooth = True
        graph1.gridlines = False

        # adding title to graph
        graph1.has_title = True
        graph_name = graph1.chart_title.text_frame
        # graph1.value_axis.minor_unit = 10
        graph_name.text = "Overall_dev_wt_weekly_performance"

        # adding x and y axis names
        category_axis_title = graph1.category_axis.axis_title
        category_axis_title.text_frame.text = "Week"
        value_axis_title = graph1.value_axis.axis_title
        value_axis_title.text_frame.text = "Score"

        Right = Inches(1)
        Top = Inches(3)
        mainWidth = Inches(3.5)
        mainHeight = Inches(2)
        mainTx_box = slide6.shapes.add_textbox(Right, Top, mainWidth, mainHeight)
        tf = mainTx_box.text_frame
        p = tf.add_paragraph()
        p.text = "DE TEAM WEEKLY\nPerformace Report"
        p.font.bold = True
        p.font.size = Pt(37)

        prs.save("slide2.pptx")

        return Response(data=employeeRating)

